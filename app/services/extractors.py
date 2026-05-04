from __future__ import annotations

import csv
from html.parser import HTMLParser
from io import StringIO
from pathlib import Path


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".csv", ".txt", ".md", ".html", ".htm", ".pptx", ".xlsx"}
GOOGLE_WORKSPACE_SHORTCUTS = {".gdoc", ".gslides", ".gsheet"}


class ExtractionError(RuntimeError):
    """Raised when a supported file cannot be converted to text."""


def is_supported_file(filename: str | Path) -> bool:
    return Path(filename).suffix.lower() in SUPPORTED_EXTENSIONS


def extract_text_from_path(path: str | Path) -> str:
    source_path = Path(path)
    extension = source_path.suffix.lower()

    if extension in GOOGLE_WORKSPACE_SHORTCUTS:
        raise ExtractionError(
            f"{source_path.name} is a Google Workspace shortcut, not the document content. "
            "Export it from Google Drive as PDF, DOCX, PPTX, XLSX, Markdown, or plain text before ingestion."
        )

    if extension not in SUPPORTED_EXTENSIONS:
        raise ExtractionError(f"Unsupported file type: {extension or 'unknown'}")

    try:
        if extension == ".pdf":
            return _extract_pdf(source_path)
        if extension == ".docx":
            return _extract_docx(source_path)
        if extension == ".csv":
            return _extract_csv(source_path)
        if extension in {".html", ".htm"}:
            return _extract_html(source_path)
        if extension == ".pptx":
            return _extract_pptx(source_path)
        if extension == ".xlsx":
            return _extract_xlsx(source_path)
        return _read_text(source_path)
    except ExtractionError:
        raise
    except Exception as exc:  # pragma: no cover - defensive wrapper
        raise ExtractionError(f"Failed to extract text from {source_path.name}: {exc}") from exc


def _read_text(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    raise ExtractionError(f"Unable to decode text file: {path.name}")


def _extract_csv(path: Path) -> str:
    raw_text = _read_text(path)
    reader = csv.reader(StringIO(raw_text))
    lines: list[str] = []

    for row in reader:
        clean_cells = [cell.strip() for cell in row]
        lines.append(", ".join(clean_cells))

    return "\n".join(lines)


class _HTMLTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._chunks: list[str] = []
        self._skip_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:  # noqa: ARG002
        if tag in {"script", "style", "noscript"}:
            self._skip_depth += 1
        if tag in {"p", "div", "section", "article", "li", "tr", "h1", "h2", "h3", "h4", "br"}:
            self._chunks.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style", "noscript"} and self._skip_depth:
            self._skip_depth -= 1
        if tag in {"p", "div", "section", "article", "li", "tr", "h1", "h2", "h3", "h4"}:
            self._chunks.append("\n")

    def handle_data(self, data: str) -> None:
        if self._skip_depth:
            return
        text = " ".join(data.split())
        if text:
            self._chunks.append(f"{text} ")

    def text(self) -> str:
        lines = [" ".join(line.split()) for line in "".join(self._chunks).splitlines()]
        return "\n".join(line for line in lines if line)


def _extract_html(path: Path) -> str:
    parser = _HTMLTextExtractor()
    parser.feed(_read_text(path))
    return parser.text()


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise ExtractionError("DOCX extraction requires the python-docx package") from exc

    document = Document(path)
    chunks: list[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            chunks.append(text)

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                chunks.append(" | ".join(cells))

    return "\n\n".join(chunks)


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ExtractionError("PPTX extraction requires the python-pptx package") from exc

    presentation = Presentation(str(path))
    slides: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = [f"Slide {slide_index}"]
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = str(shape.text).strip()
                if text:
                    parts.append(text)
        slides.append("\n".join(parts))
    return "\n\n".join(slides)


def _extract_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise ExtractionError("XLSX extraction requires the openpyxl package") from exc

    workbook = load_workbook(filename=path, read_only=True, data_only=True)
    sheets: list[str] = []
    for sheet in workbook.worksheets:
        rows: list[str] = [f"Sheet: {sheet.title}"]
        for row in sheet.iter_rows(values_only=True):
            values = ["" if cell is None else str(cell).strip() for cell in row]
            if any(values):
                rows.append(" | ".join(values))
        sheets.append("\n".join(rows))
    workbook.close()
    return "\n\n".join(sheets)


def _extract_pdf(path: Path) -> str:
    errors: list[str] = []

    try:
        import fitz

        with fitz.open(path) as document:
            pages = [page.get_text("text") or "" for page in document]
        text = "\n\n".join(page.strip() for page in pages if page.strip())
        if text:
            return text
    except ImportError:
        errors.append("pymupdf unavailable")
    except Exception as exc:
        errors.append(f"pymupdf failed: {exc}")

    try:
        import pdfplumber

        with pdfplumber.open(path) as pdf:
            pages = [page.extract_text() or "" for page in pdf.pages]
        text = "\n\n".join(page.strip() for page in pages if page.strip())
        if text:
            return text
    except ImportError as exc:
        errors.append("pdfplumber unavailable")
    except Exception as exc:
        errors.append(f"pdfplumber failed: {exc}")

    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        pages = [page.extract_text() or "" for page in reader.pages]
        text = "\n\n".join(page.strip() for page in pages if page.strip())
        if text:
            return text
    except ImportError:
        errors.append("pypdf unavailable")
    except Exception as exc:
        errors.append(f"pypdf failed: {exc}")

    detail = "; ".join(errors) if errors else "no extractable text"
    raise ExtractionError(f"PDF extraction failed for {path.name}: {detail}")
